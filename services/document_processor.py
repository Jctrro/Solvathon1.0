import os
import fitz
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract


# ===============================
# 🔍 SUPPORTED EXTENSIONS
# ===============================

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".txt", ".csv"}


def is_supported(filename: str) -> bool:
    """Check if a file extension is supported for text extraction."""
    ext = os.path.splitext(filename)[1].lower()
    return ext in SUPPORTED_EXTENSIONS


# ===============================
# 📄 UNIVERSAL TEXT EXTRACTION
# ===============================

def extract_text_universal(file_path: str) -> str:
    """
    Extract raw text from any supported document type.
    Returns plain text string.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf(file_path)

    elif ext == ".docx":
        return extract_docx(file_path)

    elif ext == ".pptx":
        return extract_pptx(file_path)

    elif ext in (".png", ".jpg", ".jpeg"):
        return extract_image(file_path)

    elif ext == ".txt":
        return extract_txt(file_path)

    elif ext == ".csv":
        return extract_csv(file_path)

    return ""


# ===============================
# 📄 STRUCTURED EXTRACTION (PAGE / SECTION AWARE)
# ===============================

def extract_text_with_structure(file_path: str) -> list[dict]:
    """
    Extract text with structural metadata (page/slide/section info).
    Returns a list of dicts:
        [{"section": "page_1", "content": "..."},  ...]

    This enables document-type-aware chunking downstream.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf_structured(file_path)

    elif ext == ".docx":
        return extract_docx_structured(file_path)

    elif ext == ".pptx":
        return extract_pptx_structured(file_path)

    elif ext in (".png", ".jpg", ".jpeg"):
        text = extract_image(file_path)
        return [{"section": "image", "content": text}]

    elif ext == ".txt":
        text = extract_txt(file_path)
        return [{"section": "full_text", "content": text}]

    elif ext == ".csv":
        text = extract_csv(file_path)
        return [{"section": "csv_data", "content": text}]

    return []


# ===============================
# 🗂️ BATCH EXTRACTION
# ===============================

def extract_multiple(file_paths: list[str]) -> list[dict]:
    """
    Process multiple files at once.
    Returns a list of dicts with filename, file_type, raw text, and structured sections.
    """
    results = []

    for fp in file_paths:
        filename = os.path.basename(fp)
        ext = os.path.splitext(fp)[1].lower().lstrip(".")

        raw_text = extract_text_universal(fp)
        sections = extract_text_with_structure(fp)

        results.append({
            "file_path": fp,
            "filename": filename,
            "file_type": ext,
            "text": raw_text,
            "sections": sections
        })

    return results


# ===============================
# 📑 INDIVIDUAL EXTRACTORS
# ===============================

def extract_pdf(file_path):
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    return text


def extract_pdf_structured(file_path) -> list[dict]:
    """Extract PDF text page-by-page with page numbers."""
    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc, start=1):
        content = page.get_text().strip()
        if content:
            pages.append({
                "section": f"page_{i}",
                "content": content
            })
    return pages


def extract_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_docx_structured(file_path) -> list[dict]:
    """
    Extract DOCX text grouped by headings.
    Each heading starts a new section; body text without a heading
    is placed under 'section_intro'.
    """
    doc = Document(file_path)
    sections = []
    current_heading = "section_intro"
    current_lines = []

    for para in doc.paragraphs:
        if para.style and para.style.name.startswith("Heading"):
            # Flush previous section
            if current_lines:
                sections.append({
                    "section": current_heading,
                    "content": "\n".join(current_lines)
                })
                current_lines = []
            current_heading = para.text.strip() or current_heading
        else:
            if para.text.strip():
                current_lines.append(para.text)

    # Flush last section
    if current_lines:
        sections.append({
            "section": current_heading,
            "content": "\n".join(current_lines)
        })

    # Fallback: if no headings found, return as one section
    if not sections:
        full = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        sections.append({"section": "full_text", "content": full})

    return sections


def extract_pptx(file_path):
    prs = Presentation(file_path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


def extract_pptx_structured(file_path) -> list[dict]:
    """Extract PPTX text slide-by-slide."""
    prs = Presentation(file_path)
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        if parts:
            slides.append({
                "section": f"slide_{i}",
                "content": "\n".join(parts)
            })

    return slides


def extract_image(file_path):
    img = Image.open(file_path)
    return pytesseract.image_to_string(img)


def extract_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_csv(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()